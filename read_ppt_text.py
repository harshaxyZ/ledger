import collections 
import collections.abc
from pptx import Presentation

prs = Presentation("ppt template.pptx")
for i, slide in enumerate(prs.slides):
    print(f"--- Slide {i+1} ---")
    for j, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            print(f"Shape {j}: {shape.text.strip()[:100]}")
