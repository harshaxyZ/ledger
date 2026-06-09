from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, content, slide_layout=1):
    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout])
    title_shape = slide.shapes.title
    title_shape.text = title
    if content:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = content
    return slide

# 1. Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Ledger: Privacy-First Offline Expense Tracker"
subtitle.text = "First-Year CSE Project Presentation\nTeam Members:\nMember 1 (USN) | Member 2 (USN) | Member 3 (USN)\nMember 4 (USN) | Member 5 (USN) | Member 6 (USN)\n\nDepartment of Computer Science and Engineering\nCollege Name\nGuide: Guide Name"

# 2. Introduction
add_slide("2. Introduction", 
"• Background: Personal finance tracking is increasingly digitized, yet many users are skeptical about data privacy and cloud dependency.\n"
"• Need: A fast, offline-capable expense tracker that requires zero cloud connectivity or user authentication.\n"
"• Significance: Encourages financial literacy by lowering the barrier to entry while ensuring zero data harvesting and complete user control.")

# 3. Problem Statement
add_slide("3. Problem Statement",
"• Description: Most modern expense trackers (e.g., Mint, YNAB) force users into cloud ecosystems, exposing sensitive financial data to corporate servers, advertisers, or breaches.\n"
"• Motivation: To build a solution that entirely operates on the local device, solving the privacy flaw while maintaining modern features like analytics, gamification, and insights locally.")

# 4. Objectives
add_slide("4. Objectives",
"• Objective 1: Develop an offline-first PWA (Progressive Web App) for tracking daily income and expenses without any backend database.\n"
"• Objective 2: Implement smart budgeting with real-time warnings and visual dashboards entirely via LocalStorage.\n"
"• Objective 3: Integrate gamification (streaks) and privacy-centric coaching without storing data externally.\n"
"• Expected Outcome: A robust, zero-latency financial dashboard that rivals premium apps but respects user privacy.")

# 5. Literature Survey 1
slide = add_slide("5. Literature Survey", "")
rows, cols = 4, 4
left, top, width, height = Inches(0.5), Inches(1.5), Inches(9.0), Inches(2.5)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table
table.columns[0].width = Inches(1.5)
table.columns[1].width = Inches(2.5)
table.columns[2].width = Inches(2.5)
table.columns[3].width = Inches(2.5)
headers = ['Author/Year', 'Existing Work', 'Findings', 'Limitations']
for i, h in enumerate(headers): table.cell(0, i).text = h
table.cell(1, 0).text = "Intuit (2020)"
table.cell(1, 1).text = "Mint / Cloud Trackers"
table.cell(1, 2).text = "Aggregated cloud data"
table.cell(1, 3).text = "Data selling, mandatory internet"
table.cell(2, 0).text = "YNAB (2022)"
table.cell(2, 1).text = "Zero-based budgeting"
table.cell(2, 2).text = "Effective saving strategy"
table.cell(2, 3).text = "High subscription cost, cloud sync"
table.cell(3, 0).text = "Expensify (2021)"
table.cell(3, 1).text = "Corporate expense app"
table.cell(3, 2).text = "Automated receipt scanning"
table.cell(3, 3).text = "Bloated for personal use"

# 6. Existing System
add_slide("6. Existing System",
"• Current Solutions: Mint, YNAB, Spendee, Wallet.\n"
"• Advantages: Cross-device sync, bank account integration, community support.\n"
"• Limitations:\n"
"  - Requires mandatory account creation.\n"
"  - Financial data is stored on remote servers.\n"
"  - Not usable without an internet connection.\n"
"  - Prone to latency, loading screens, and server downtimes.")

# 7. Proposed System
add_slide("7. Proposed System",
"• Overview: 'Ledger' is a purely client-side React PWA that handles all logic within the browser.\n"
"• Features:\n"
"  - Instant transactions, offline history.\n"
"  - Spending insights, sparkline trends, and category analysis.\n"
"  - Streak gamification to encourage daily logging.\n"
"  - Safe local data export (JSON backup).\n"
"• Benefits: Absolute privacy, zero loading times, free forever, installable on mobile via PWA.")

# 8. Methodology
add_slide("8. Methodology",
"• Step 1: Requirement Analysis - Identify core features (Tracking, Budgets, AI, PWA).\n"
"• Step 2: UI/UX Design - Create a sleek, dark-themed, mobile-responsive layout.\n"
"• Step 3: Frontend Setup - Initialize Vite + React and configure TailwindCSS.\n"
"• Step 4: Logic Implementation - Build custom React hooks to manage LocalStorage seamlessly.\n"
"• Step 5: PWA Integration - Configure service workers for offline caching and installation.\n"
"• Step 6: Testing & QA - Verify standalone mode behaviors, budget warnings, and edge cases.")

# 9. System Architecture (1/2)
add_slide("9. System Architecture (1/2)",
"• Client-Side Only Architecture:\n"
"  - The app lives entirely on the client's device.\n"
"  - Browser LocalStorage acts as the persistent database.\n"
"  - React handles UI updates instantly without network latency.\n"
"  - Service Workers intercept requests to serve cached files offline.")

# 10. System Architecture (2/2)
add_slide("10. System Architecture (2/2)",
"• Block Diagram Flow:\n"
"  [User Input / UI Interaction]\n"
"          ↓\n"
"  [React Custom Hooks (useTransactions, useBudgets)]\n"
"          ↓\n"
"  [State Update & LocalStorage Sync]\n"
"          ↓\n"
"  [UI Re-render (Charts, History, Budgets)]")

# 11. Technologies Used
add_slide("11. Technologies Used",
"• Programming Languages:\n"
"  - JavaScript (ES6+), HTML5, CSS3.\n"
"• Development Tools/Software:\n"
"  - React (UI Library).\n"
"  - TailwindCSS (Utility-first styling).\n"
"  - Vite (Build Tool & Dev Server).\n"
"  - Framer Motion (Animations & Page Transitions).\n"
"• Database / Storage:\n"
"  - Browser LocalStorage API (No external database).")

# 12. Implementation
add_slide("12. Implementation",
"• Modules Developed:\n"
"  1. Transaction Engine: Add, delete, and categorize income/expenses.\n"
"  2. Budget Watchdog: Real-time calculation of category limits vs. spending.\n"
"  3. Visual Analytics: Doughnut charts and trend sparklines for quick insights.\n"
"  4. AI Advisor: Context-aware financial tips.\n"
"  5. PWA Wrapper: Offline detection, 'Add to Homescreen' prompts, sync handling.")

# 13. Results and Discussion (1/2)
add_slide("13. Results and Discussion (1/2)",
"• Outputs Obtained:\n"
"  - A fully functional, zero-latency Progressive Web App.\n"
"  - Responsive design that looks native on both Mobile and Desktop.\n"
"  - Instant gamification updates (streaks) based on local timestamps.")

# 14. Results and Discussion (2/2)
add_slide("14. Results and Discussion (2/2)",
"• UI Walkthrough (Screenshots):\n"
"  - Dashboard: Highlights net balance, sparkline trend, and daily insights.\n"
"  - History: Pill-button filtering with beautiful scrolling margins.\n"
"  - Insights: Dynamic doughnut charts excluding income for accurate tracking.\n"
"  - Settings: Budget thresholds and local data export controls.")

# 15. Conclusion
add_slide("15. Conclusion",
"• Summary: Ledger successfully reimagines the personal finance tracker by removing the cloud dependency entirely.\n"
"• Objectives Achieved: Built an offline-first, aesthetically premium, privacy-respecting dashboard.\n"
"• Key Learnings: Mastered React state management, PWA lifecycle hooks, responsive design with Tailwind, and building robust client-side storage architectures.")

# 16. Future Enhancements
add_slide("16. Future Enhancements",
"• End-to-End Encryption: Implement local AES-256 encryption before saving to LocalStorage.\n"
"• Multi-Device Sync: Allow secure, peer-to-peer WebRTC syncing between a user's phone and laptop without a central server.\n"
"• Custom Categories: Let users create their own spending categories and assign custom colors/icons.\n"
"• Receipt Scanning: Integrate a local, lightweight OCR library (e.g., Tesseract.js) to scan bills on-device.")

# 17. References
add_slide("17. References",
"• React Documentation: https://react.dev/\n"
"• Vite PWA Plugin Docs: https://vite-pwa-org.netlify.app/\n"
"• TailwindCSS Framework: https://tailwindcss.com/\n"
"• Framer Motion Library: https://www.framer.com/motion/\n"
"• MDN Web Docs (LocalStorage & Service Workers): https://developer.mozilla.org/")

# 18. Thank You
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Thank You!"
subtitle.text = "Any Questions?"

prs.save("PPT.pptx")
print("PPT Generated Successfully!")
