import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import docx
from docx.shared import Pt as DocxPt
import os

ppt_path = "FIRST SLIDE.pptx"
prs = Presentation(ppt_path)

# Ensure dark cinematic theme
BG_COLOR = RGBColor(15, 15, 18)
ACCENT_COLOR = RGBColor(220, 20, 60)
TEXT_COLOR = RGBColor(240, 240, 240)
SUBTEXT_COLOR = RGBColor(160, 160, 160)

def add_slide(title, content_list, subtitle=""):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title.upper()
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = TEXT_COLOR
    p.font.name = "Arial"
    
    # Subtitle or Red Accent Line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(1.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8), Inches(0.5))
        sp = sub_box.text_frame.add_paragraph()
        sp.text = subtitle
        sp.font.size = Pt(20)
        sp.font.color.rgb = ACCENT_COLOR
        sp.font.name = "Arial"
        sp.font.bold = True

    # Content
    top = Inches(2.2) if not subtitle else Inches(2.5)
    content_box = slide.shapes.add_textbox(Inches(0.8), top, Inches(8.4), Inches(4.5))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    
    for item in content_list:
        p = ctf.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = SUBTEXT_COLOR
        p.font.name = "Arial"
        p.space_after = Pt(20)
    
    return slide

# Generating the Slides
add_slide("Project Details", ["Project Name: LEDGER", "A Privacy-First Offline Expense Tracker", "", "Department of Computer Science & Engineering", "Under the Guidance of: [Your Guide's Name]"], "Team Presentation")
add_slide("Introduction", ["Modern finance apps are cluttered, invasive, and ad-heavy.", "Users are forced to surrender personal data to third-party servers.", "Ledger re-imagines expense tracking as a completely private experience.", "It lives directly on your device, ensuring zero data harvesting."])
add_slide("Need For Project", ["Data Privacy Crisis: Financial data is sold to advertisers.", "Bloatware: Most apps have confusing UI and unnecessary features.", "Internet Dependency: Tracking fails when offline or in bad network zones.", "Friction: Adding a transaction takes too many clicks."])
add_slide("Problem Statement", ["To develop a zero-latency, offline-capable progressive web application (PWA) for expense tracking that entirely eliminates backend server dependency, prioritizes strict user data privacy, and utilizes a gamified approach to encourage consistent financial habits."])
add_slide("Objectives", ["1. Achieve 100% offline functionality using local storage mechanisms.", "2. Provide instant, sub-second transaction logging.", "3. Implement localized AI insights without compromising data.", "4. Deliver a premium, distraction-free user interface."])
add_slide("Literature Survey Table", ["1. Existing App (Splitwise) | Server dependent | Cluttered UI", "2. Existing App (Mint) | Data harvesting | Shut down/Ads", "3. Existing App (Wallet) | Paid features | Complex setup", "4. LEDGER (Proposed) | 100% Offline | Zero Ads, Private"])
add_slide("Literature Survey Analysis", ["Analysis shows a massive gap for privacy-focused financial tools.", "Most alternatives monetize user data to sustain their free tier.", "Complex cloud-syncing creates latency and requires constant connectivity.", "Local-first applications are gaining traction in modern SaaS."])
add_slide("Existing System", ["Current systems rely on centralized cloud databases (Firebase, AWS).", "They require mandatory user authentication (OAuth, Email).", "They mandate constant internet connectivity for basic operations.", "Data is subject to server breaches and third-party tracking."])
add_slide("Limitations Of Existing Systems", ["• High Latency: Waiting for cloud sync on every launch.", "• Privacy Risks: Sensitive spending patterns are exposed.", "• Feature Creep: Constant push notifications for premium upgrades.", "• Vulnerability: Server downtimes render the app useless."])
add_slide("Proposed System", ["Ledger shifts all computation to the client side.", "Data is persistently stored in the browser's IndexedDB/LocalStorage.", "Service Workers cache the application core for instant offline booting.", "Groq-powered AI coaching is invoked only on demand, completely amnesic."])
add_slide("Key Features", ["• Sub-second Transaction Logging", "• Dynamic Budget Warnings", "• Categorized Spending Insights", "• Habit-forming Streak Gamification", "• 'Ledger Coach' AI Advisor", "• Native App Installation (PWA)"])
add_slide("Methodology", ["1. Requirement Analysis (Privacy & Speed)", "2. UI/UX Prototyping (Dark Cinematic Theme)", "3. Frontend Development (React + Tailwind)", "4. Local Data Architecture (Hooks & Storage)", "5. PWA Integration (Vite PWA Plugin)", "6. Testing & Deployment (Vercel)"])
add_slide("System Architecture", ["USER INTERFACE (React Components)", "          ↓", "APPLICATION LOGIC (Custom Hooks)", "          ↓", "LOCAL STORAGE (kuber_transactions)", "          ↓", "ANALYTICS ENGINE (Real-time Math)", "          ↓", "UI UPDATES (Instant Rendering)"])
add_slide("Technology Stack", ["• Core: React 19, JavaScript (ES6+)", "• Build Tool: Vite", "• Styling: Tailwind CSS v4, Framer Motion", "• Storage: Browser LocalStorage API", "• Offline Tech: vite-plugin-pwa (Workbox)", "• Deployment: Vercel"])
add_slide("Implementation Modules", ["Module 1: The Transaction Engine (Core Logging)", "Module 2: The Budget System (Limits & Warnings)", "Module 3: The Analytics Dashboard (Charts & Trends)", "Module 4: The AI Coach (Groq API Integration)", "Module 5: The PWA Layer (Manifest & Service Worker)"])
add_slide("UI Showcase", ["* Dashboard: Clean, matte black aesthetic with crimson warnings.", "* History: Detailed, searchable transaction list.", "* Coach: Chat interface that reads local context.", "* Settings: Instant data export and reset controls."])
add_slide("Results And Discussion", ["The application successfully achieves 0ms latency for logging.", "PWA installation allows it to function as a native mobile app.", "All tests confirm that no data leaves the device except for intentional AI queries.", "The UI successfully mimics a premium, million-dollar product."])
add_slide("Benefits And Impact", ["• Peace of Mind: Users know their finances are strictly private.", "• Accessibility: Works anywhere, even completely off-grid.", "• Financial Literacy: AI Coach provides tailored, unbiased advice.", "• Productivity: Removing friction leads to better habit formation."])
add_slide("Future Enhancements", ["• Custom User-Defined Categories", "• Recurring Transaction Automation", "• Biometric Authentication (FaceID/Fingerprint)", "• Cross-Device Peer-to-Peer Sync (E2E Encrypted)"])
add_slide("Conclusion", ["Ledger proves that powerful financial tracking does not require sacrificing privacy.", "By leveraging modern web capabilities, we delivered a native-feeling application.", "It is fast, secure, and ready to scale entirely on the client side."])
add_slide("References", ["1. React Documentation (react.dev)", "2. Tailwind CSS Documentation", "3. Progressive Web Apps (MDN Web Docs)", "4. LocalStorage & IndexedDB API Specs", "5. Groq AI API Guidelines"])
add_slide("Thank You", ["LEDGER", "Private by Design.", "", "Questions?"])

prs.save("LEDGER_Final_Presentation.pptx")

# Now generate Script DOCX
doc = docx.Document()
doc.add_heading('Ledger Presentation Script', 0)

script_data = [
    ("Presenter 1 (Slides 2-4)", [
        ("Slide 2: Project Details", "Welcome everyone. We are thrilled to present Ledger—a privacy-first, offline expense tracker. Let's dive in."),
        ("Slide 3: Introduction", "Most finance apps today are invasive. They force you to create accounts and hand over your data. Ledger is different. It lives entirely on your device, ensuring zero data harvesting."),
        ("Slide 4: Need For Project", "The financial app space is plagued by privacy risks and bloatware. Furthermore, many of them fail to work offline, adding friction to simply logging a coffee purchase.")
    ]),
    ("Presenter 2 (Slides 5-7)", [
        ("Slide 5: Problem Statement", "Our goal was simple: To develop a zero-latency, offline-capable PWA that completely eliminates backend dependency, prioritizes your privacy, and builds positive financial habits."),
        ("Slide 6: Objectives", "We had 4 objectives: Achieve 100% offline functionality, provide instant sub-second logging, implement smart local insights, and deliver a premium, distraction-free UI."),
        ("Slide 7: Literature Survey Table", "Looking at existing systems like Splitwise or Mint, they either rely heavily on servers, harvest your data, or lock basic features behind paywalls. Ledger does none of that.")
    ]),
    ("Presenter 3 (Slides 8-10)", [
        ("Slide 8: Literature Survey Analysis", "Our analysis confirmed a massive gap. The industry standard is to monetize user data. By shifting to a local-first approach, we found we could deliver a vastly superior user experience."),
        ("Slide 9: Existing System", "Current systems use centralized cloud databases and mandatory authentication. This means you need constant internet, and your data is subject to breaches."),
        ("Slide 10: Limitations", "These legacy systems suffer from high latency during cloud syncs, severe privacy risks, and annoying feature creep pushing you to upgrade.")
    ]),
    ("Presenter 4 (Slides 11-13)", [
        ("Slide 11: Proposed System", "Our proposed system, Ledger, shifts all computation to the client side. Data is persistently stored in your browser, and Service Workers cache the app for instant offline booting."),
        ("Slide 12: Key Features", "The standout features include sub-second logging, dynamic budget warnings that turn crimson red when you overspend, habit-forming streaks, and our intelligent AI Coach."),
        ("Slide 13: Methodology", "We followed a streamlined methodology: starting with strict privacy requirements, moving to a dark cinematic UI/UX prototype, developing the frontend in React, and finalizing with PWA integration.")
    ]),
    ("Presenter 5 (Slides 14-16)", [
        ("Slide 14: System Architecture", "Our architecture is elegantly simple. The React UI triggers custom hooks, which write directly to LocalStorage. Our Analytics Engine then instantly reads this and updates the screen. No servers required."),
        ("Slide 15: Technology Stack", "We built this using React 19 and Vite for the core, Tailwind CSS for the premium styling, and the browser's native LocalStorage and Service Worker APIs for the offline capabilities."),
        ("Slide 16: Implementation Modules", "The build was divided into 5 modules: The Transaction Engine, the Budget System, the Analytics Dashboard, the AI Coach integration via Groq, and the PWA layer.")
    ]),
    ("Presenter 6 (Slides 17-23)", [
        ("Slide 17: UI Showcase", "As you can see from our app, the dashboard is clean and matte black. The history is highly searchable, the Coach reads your local context, and you can export your data anytime."),
        ("Slide 18: Results And Discussion", "The result is exactly what we aimed for: 0ms latency, native mobile feel via PWA installation, and 100% data privacy. It feels like a premium product."),
        ("Slide 19: Benefits And Impact", "The true impact is peace of mind. Users can track their money completely off-grid, and our AI coach provides unbiased financial literacy without storing your history."),
        ("Slide 20: Future Enhancements", "In the future, we plan to add custom categories, recurring automated transactions, and biometric locks like FaceID for extra local security."),
        ("Slide 21: Conclusion", "In conclusion, Ledger proves that powerful financial tracking doesn't require sacrificing privacy. We've built a fast, secure, and highly scalable tool."),
        ("Slide 22: References", "Our development relied on React Docs, Tailwind Docs, and MDN Web Specs for PWA architecture."),
        ("Slide 23: Thank You", "Thank you for your time. Ledger is private by design. We're now open for any questions.")
    ])
]

for presenter, lines in script_data:
    doc.add_heading(presenter, level=1)
    for title, text in lines:
        p = doc.add_paragraph()
        p.add_run(title + ": ").bold = True
        p.add_run(text)

doc.save('Ledger_Presentation_Script.docx')
print("Successfully generated files!")
