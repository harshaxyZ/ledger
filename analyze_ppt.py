import collections 
import collections.abc
from pptx import Presentation

prs = Presentation("ppt template.pptx")
print(f"Total slides: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    shapes = [s for s in slide.shapes if s.has_text_frame]
    print(f"Slide {i+1}: {len(shapes)} text shapes")
